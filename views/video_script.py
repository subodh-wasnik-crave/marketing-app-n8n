import streamlit as st
import requests
import json
from PyPDF2 import PdfReader  # type: ignore
from docx import Document as DocxDocument  # type: ignore
from pptx import Presentation  # type: ignore
import markdown

def show(navigate_to):

    col_nav, _ = st.columns([1, 5])
    with col_nav:
        if st.button("← Home"):
            navigate_to("Home")
    
    # -----------------------------------------------------------------------------
    # CONFIGURATION
    # -----------------------------------------------------------------------------
    st.set_page_config(layout="wide", page_title="Video Script Generator")

    # REPLACE THIS WITH YOUR ACTUAL N8N WEBHOOK URL
    # N8N_WEBHOOK_URL = "http://localhost:5678/webhook/Video-script"
    N8N_WEBHOOK_URL = st.secrets.get("n8n").get("video_script_api")

    # Initialize Session State
    if "video_output" not in st.session_state:
        st.session_state.video_output = ""

    if "last_params" not in st.session_state:
        st.session_state.last_params = {}  # Stores context for refinement

    # -----------------------------------------------------------------------------
    # HELPER FUNCTIONS
    # -----------------------------------------------------------------------------
    def extract_text_from_file(file):
        text = ""
        name = file.name.lower()
        try:
            if name.endswith(".txt"):
                text = file.read().decode("utf-8", errors="ignore")
            elif name.endswith(".pdf"):
                pdf = PdfReader(file)
                for page in pdf.pages:
                    text += page.extract_text() or ""
            elif name.endswith(".docx"):
                doc = DocxDocument(file)
                text = "\n".join([p.text for p in doc.paragraphs])
            elif name.endswith(".pptx"):
                ppt = Presentation(file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        except Exception:
            pass
        return text.strip()

    # -----------------------------------------------------------------------------
    # SIDEBAR: CONTENT CONFIGURATION
    # -----------------------------------------------------------------------------
    with st.sidebar:
        st.markdown("## ⚙️ Content Configuration")

        tone = st.selectbox(
            "🎨 Tone",
            [
                "Professional", "Friendly", "Authoritative", "Playful", "Inspirational",
                "Conversational", "Casual", "Semi-casual", "Business professional",
                "Approachable", "Informative", "Assertive", "Engaging",
                "Visionary (for Thought Leadership)", "Confident", "Data-driven",
                "Plainspoken / Direct", "Witty", "Storytelling"
            ],
            key="video_tone",
        )

        target_audience = st.selectbox(
            "🎯 Target Audience",
            ["Senior Management", "Middle Management", "Junior / Entry Level Staff"],
            key="video_audience",
        )

        industry = st.text_input(
            "🏢 Industry (optional)",
            placeholder="e.g., Manufacturing, Retail, Technology",
            key="video_industry",
        )

        time_limit = st.slider("⏱️ Video Duration (minutes)", 0.5, 10.0, 1.5, step=0.5, key="time_limit")

        cta_options = [
            "Talk to our experts", "Learn more about our solutions", "Book a free consultation",
            "Book Assessment", "Contact us today", "Download the full guide", "Request a demo",
        ]
        cta_choice = st.selectbox("📢 Call-to-Action (CTA)", cta_options, key="video_cta_choice")
        

    # -----------------------------------------------------------------------------
    # TOP SECTION: FILES & URLS
    # -----------------------------------------------------------------------------
    st.title("🚀 Video Script Generator")

    col1, col2 = st.columns([1.2, 1.8])

    with col1:
        st.markdown("#### 📎 Upload Documents")
        uploaded_files = st.file_uploader(
            "Upload Reference Files (TXT, PDF, DOCX, PPTX)",
            type=["txt", "pdf", "docx", "pptx"],
            accept_multiple_files=True,
            label_visibility="collapsed",
            key="video_files"
        )
        
        # Process files immediately to be ready for generation
        file_context = ""
        if uploaded_files:
            for f in uploaded_files:
                extracted_text = extract_text_from_file(f) 
                file_context += f"--- Content from {f.name} ---\n{extracted_text}\n\n"

    with col2:
        st.markdown("#### 🔗 Reference URLs")
        reference_urls = st.text_area(
            "Add Reference URLs (comma-separated)",
            placeholder="https://example.com/page1, https://example.com/page2",
            height=70,
            label_visibility="collapsed",
            key="video_urls"
        )
        url_list = [url.strip() for url in reference_urls.split(",") if url.strip()]

    st.divider()

    # -----------------------------------------------------------------------------
    # MAIN LAYOUT: INPUTS VS OUTPUT
    # -----------------------------------------------------------------------------
    left, right = st.columns([1, 2])

    with left:
        st.markdown("### ✍️ Topic")
        query = st.text_input("Video Topic", key="video_query")
        

        st.markdown("<br>", unsafe_allow_html=True)
        generate_button = st.button("Generate Video", type="primary", use_container_width=True)
        # -------------------------------------------------------------------------
        # REFINE SECTION (Only shows if we have output)
        # -------------------------------------------------------------------------
        if st.session_state.video_output:
            st.markdown("---")
            st.markdown("### 🛠️ Refine Content")
            # st.info("Ask for changes (e.g., 'Make it shorter', 'Add more statistics'). Context is preserved.")
            
            refine_instruction = st.text_area(
                "Refinement Instruction:",
                height=80,
                placeholder="What should be changed?",
                key="video_refine_input"
            )
            
            apply_refine = st.button("Apply Changes", use_container_width=True)
            

    with right:
        # st.markdown("### 📝 Output")
        
        # -------------------------------------------------------------------------
        # LOGIC: GENERATE NEW video
        # -------------------------------------------------------------------------
        if generate_button and query:
            with st.spinner("🚀 Generating video via n8n..."):
                
                # 1. CAPTURE CONTEXT
                # We save all inputs to session_state so we can re-send them during refinement
                st.session_state.last_params = {
                    "query": query,
                    "tone": tone,
                    "target_audience": target_audience,
                    "industry": industry,
                    "time_limit": time_limit,
                    "cta_choice": cta_choice,
                    "reference_urls": url_list,
                    "reference_file_content": file_context  # CRITICAL: Keeps file text in memory
                }

                # 2. PREPARE PAYLOAD
                payload = {
                    "action": "generate",
                    **st.session_state.last_params # Unpack all params
                }

                try:
                    response = requests.post(N8N_WEBHOOK_URL, json=payload)
                    
                    if response.status_code == 200:
                        try:
                            data = response.json()
                            result_text = data.get("output", data.get("text", str(data)))
                        except:
                            result_text = response.text
                            
                        st.session_state.video_output = result_text
                        st.success("video generated successfully!")
                        st.rerun() # Rerun to show the Refine options
                    else:
                        st.error(f"Error {response.status_code}: {response.text}")
                        
                except Exception as e:
                    st.error(f"Connection Error: {e}")

        # -------------------------------------------------------------------------
        # LOGIC: REFINE EXISTING video
        # -------------------------------------------------------------------------
        # We check if 'apply_refine' exists because it's conditionally rendered above
        if 'apply_refine' in locals() and apply_refine and st.session_state.video_output:
            with st.spinner("✨ Refining content (preserving context)..."):
                
                # 1. RETRIEVE CONTEXT
                # We fetch the params used during the INITIAL generation
                context_params = st.session_state.last_params
                
                # 2. PREPARE PAYLOAD
                # We mix the old context with the new instruction and current text
                payload = {
                    "action": "refine",
                    "current_video_content": st.session_state.video_output,
                    "refine_instruction": refine_instruction,
                    
                    # Re-send original context so n8n can access it again
                    "query": context_params.get("query", ""),
                    "tone": context_params.get("tone", "Professional"),
                    "target_audience": context_params.get("target_audience", "Senior Management"),
                    "reference_file_content": context_params.get("reference_file_content", ""),
                    "reference_urls": context_params.get("reference_urls", []),
                    "primary_keyword": context_params.get("primary_keyword", ""),
                    "time_limit": context_params.get("time_limit", 1.5)
                }
                
                try:
                    response = requests.post(N8N_WEBHOOK_URL, json=payload)
                    
                    if response.status_code == 200:
                        try:
                            data = response.json()
                            result_text = data.get("output", data.get("text", str(data)))
                        except:
                            result_text = response.text
                            
                        st.session_state.video_output = result_text
                        st.success("Refinement applied!")
                        st.rerun()
                    else:
                        st.error("Failed to refine content.")
                except Exception as e:
                    st.error(f"Error: {e}")


        # -------------------------------------------------------------------------
        # DISPLAY OUTPUT CONTAINER
        # -------------------------------------------------------------------------
        if st.session_state.video_output:
            container = st.container(border=True)
            with container:
                st.markdown(st.session_state.video_output)
                
            st.download_button(
                label="Download Markdown",
                data=st.session_state.video_output,
                file_name=f"video_{st.session_state.last_params.get('query', 'output')[:10].replace(' ', '_')}.md",
                mime="text/markdown"
            )
