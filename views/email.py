import streamlit as st
import pandas as pd
import requests
import json
import ast  # Added for parsing dictionary strings
import re   # Added for regex fallback
from io import StringIO
from PyPDF2 import PdfReader  # type: ignore
from docx import Document as DocxDocument  # type: ignore
from pptx import Presentation  # type: ignore


def show(navigate_to):


    col_nav, _ = st.columns([1, 5])
    with col_nav:
        if st.button("← Home"):
            navigate_to("Home")
    
    
    # --- CONFIGURATION ---
    # GENERATE_WEBHOOK_URL = "http://localhost:5678/webhook/customised-email-generation"
    # REFINE_WEBHOOK_URL = "http://localhost:5678/webhook/refine-email"
    # SEND_WEBHOOK_URL = "http://localhost:5678/webhook/send-email" 
    
    GENERATE_WEBHOOK_URL = st.secrets.get("n8n").get("email_generate_api")
    REFINE_WEBHOOK_URL = st.secrets.get("n8n").get("email_refine_api")

    st.set_page_config(page_title="Email Generator", layout="wide")

    # --- SESSION STATE INITIALIZATION ---
    if "leads_df" not in st.session_state:
        st.session_state.leads_df = None
    if "processing" not in st.session_state:
        st.session_state.processing = False
    if "file_context" not in st.session_state:
        st.session_state.file_context = ""

    # --- HELPER FUNCTIONS ---
    def extract_text_from_file(file):
        text = ""
        name = file.name.lower()
        try:
            if name.endswith(".txt"):
                text = file.read().decode("utf-8", errors="ignore")
            elif name.endswith(".pdf"):
                pdf = PdfReader(file)
                for page in pdf.pages:
                    text += page.extract_text() or ""
            elif name.endswith(".docx"):
                doc = DocxDocument(file)
                text = "\n".join([p.text for p in doc.paragraphs])
            elif name.endswith(".pptx"):
                ppt = Presentation(file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        except Exception:
            pass
        return text.strip()

    def parse_llm_response(content):
        """
        Robust parsing function that handles:
        1. Valid JSON
        2. Python Dictionary Strings (ast.literal_eval)
        3. Dirty/Malformed strings using Regex
        """
        if not isinstance(content, str):
            if isinstance(content, dict):
                return content
            return {"body": str(content)}

        parsed = None
        
        # Attempt 1: Standard JSON
        try:
            parsed = json.loads(content)
            if isinstance(parsed, dict): return parsed
        except: pass

        # Attempt 2: AST Literal Eval (for Python dict strings with single quotes)
        try:
            # Sanitize raw newlines which break ast.literal_eval
            # We replace actual newlines with literal \n characters for the eval
            sanitized_content = content.replace('\n', '\\n') 
            parsed = ast.literal_eval(sanitized_content)
            if isinstance(parsed, dict): return parsed
        except: pass
        
        # Attempt 3: Regex Extraction (The "Nuclear Option")
        # This looks for 'subject': '...' or "subject": "..." patterns
        try:
            subject_match = re.search(r"['\"]subject['\"]\s*:\s*['\"](.*?)['\"](?:\s*,|\s*})", content, re.IGNORECASE | re.DOTALL)
            body_match = re.search(r"['\"]body['\"]\s*:\s*['\"](.*?)['\"](?:\s*,|\s*})", content, re.IGNORECASE | re.DOTALL)
            
            if subject_match or body_match:
                return {
                    "subject": subject_match.group(1) if subject_match else "",
                    "body": body_match.group(1) if body_match else ""
                }
        except: pass
        
        # Return as raw text if all else fails
        return {"subject": "", "body": content}

    def generate_bulk_emails(df, common_params):
        progress_bar = st.progress(0)
        total_rows = len(df)
        generated_emails = []
        
        for index, row in df.iterrows():
            payload = {
                "action": "generate",
                "first_name": row.get('first_name', ''),
                "last_name": row.get('last_name', ''),
                "email": row.get('email', ''),
                "org_name": row.get('org_name', ''),
                **common_params
            }
            
            try:
                response = requests.post(GENERATE_WEBHOOK_URL, json=payload)
                if response.status_code == 200:
                    data = response.json()
                    
                    # Check if data itself is the dict we want
                    if isinstance(data, dict) and "subject" in data and "body" in data:
                        clean_data = data
                    else:
                        # If nested inside 'output' or 'text'
                        raw_text = data.get("output", data.get("text", str(data)))
                        # Try to parse this inner text immediately
                        clean_data = parse_llm_response(raw_text)
                    
                    # Normalize keys
                    clean_data = {k.lower(): v for k, v in clean_data.items()}
                    
                    # Store as Valid JSON String
                    email_content = json.dumps({
                        "subject": clean_data.get("subject", ""),
                        "body": clean_data.get("body", clean_data.get("text", ""))
                    }, ensure_ascii=False)
                    
                else:
                    email_content = json.dumps({"subject": "Error", "body": f"Error: {response.status_code}"})
            except Exception as e:
                email_content = json.dumps({"subject": "Connection Error", "body": str(e)})
                
            generated_emails.append(email_content)
            progress_bar.progress((index + 1) / total_rows)
            
        df['Generated Email'] = generated_emails
        df['Status'] = 'Draft' 
        return df

    def send_final_emails(df):
        count = 0
        progress_bar = st.progress(0)
        total = len(df)
        
        for index, row in df.iterrows():
            if row['Status'] in ['Draft', 'Refined', 'Approved']:
                
                raw_content = row['Generated Email']
                parsed = parse_llm_response(raw_content)
                
                payload = {
                    "email": row['email'],
                    "subject": parsed.get('subject', "Your Custom Subject"),
                    "body": parsed.get('body', raw_content)
                }
                try:
                    requests.post(SEND_WEBHOOK_URL, json=payload)
                    st.session_state.leads_df.at[index, 'Status'] = 'Sent'
                    count += 1
                except:
                    st.session_state.leads_df.at[index, 'Status'] = 'Failed'
            progress_bar.progress((index + 1) / total)
            
        return count

    @st.dialog("Review & Refine Email", width="large")
    def email_editor_dialog(index, row_data, common_params):
        st.write(f"**Lead:** {row_data.get('first_name')} {row_data.get('last_name')} | **Org:** {row_data.get('org_name')}")
        
        # 1. PARSE CONTENT
        raw_content = row_data['Generated Email']
        parsed_data = parse_llm_response(raw_content)

        # 2. Extract specific fields
        current_subject = parsed_data.get("subject", "")
        current_body = parsed_data.get("body", "")

        # Cleanup: If body still contains escaped newlines literal characters, fix them for display
        # (e.g., if it shows "Hi\n\nThere", we want actual line breaks)
        if isinstance(current_body, str):
            current_body = current_body.replace("\\n", "\n")

        col1, col2 = st.columns([2, 1])
        
        with col1:
            new_subject = st.text_input("Subject Line", value=current_subject)
            new_body = st.text_area("Email Body", value=current_body, height=400)
        
        with col2:
            st.subheader("Refine with AI")
            refine_instruction = st.text_area("Instructions", placeholder="Make it shorter, more professional...")
            
            if st.button("✨ Refine Email", type="primary"):
                with st.spinner("Refining..."):
                    payload = {
                        "current_subject": new_subject,
                        "current_email": new_body,
                        "instruction": refine_instruction,
                        **common_params 
                    }
                    try:
                        response = requests.post(REFINE_WEBHOOK_URL, json=payload)
                        if response.status_code == 200:
                            data = response.json()
                            # Handle list wrapper
                            if isinstance(data, list): data = data[0] if len(data) > 0 else {}
                            
                            # Handle if response is just the body text or a dict
                            if isinstance(data, dict):
                                refined_body_text = data.get("output", data.get("text", str(data)))
                            else:
                                refined_body_text = str(data)

                            # Clean response
                            refined_parsed = parse_llm_response(refined_body_text)
                            # If the refined response is just text, use it as body. If it's json, parse it.
                            final_body = refined_parsed.get("body") if refined_parsed.get("body") else refined_body_text

                            updated_full_content = json.dumps({
                                "subject": new_subject,
                                "body": final_body
                            }, ensure_ascii=False)
                            
                            st.session_state.leads_df.at[index, 'Generated Email'] = updated_full_content
                            st.session_state.leads_df.at[index, 'Status'] = 'Refined'
                            st.success("Refined!")
                            st.rerun()
                        else:
                            st.error("Refinement failed.")
                    except Exception as e:
                        st.error(f"Error: {e}")

        if st.button("✅ Approve & Save"):
            final_content = json.dumps({
                "subject": new_subject,
                "body": new_body
            }, ensure_ascii=False)
            st.session_state.leads_df.at[index, 'Generated Email'] = final_content
            st.session_state.leads_df.at[index, 'Status'] = 'Approved'
            st.rerun()

    # --- SIDEBAR UI ---
    with st.sidebar:
        st.markdown("## ⚙️ Content Configuration")
        tone = st.selectbox(
            "🎨 Tone",
            [
                "Professional", "Friendly", "Authoritative", "Playful", "Inspirational",
                "Conversational", "Casual", "Semi-casual", "Business professional",
                "Approachable", "Informative", "Assertive", "Engaging",
                "Visionary (for Thought Leadership)", "Confident", "Data-driven",
                "Plainspoken / Direct", "Witty", "Storytelling"
            ],
            key="video_tone",
        )

        word_limit = st.slider("📝 Word Limit", 300, 2000, 500, step=100, key="email_word_limit")

        cta_options = [
            "Talk to our experts", "Learn more about our solutions", "Book a free consultation",
            "Book Assessment", "Contact us today", "Download the full guide", "Request a demo",
        ]
        cta_choice = st.selectbox("📢 Call-to-Action (CTA)", cta_options, key="email_cta_choice")
        
        st.divider()
        
        st.header("Upload Leads")
        st.write("Upload Lead CSV with columns: first_name, last_name, email, org_name")
        uploaded_file = st.file_uploader(
            "Choose a CSV file", 
            type=["csv", "excel"],
            key="lead_file"
        )
        
        if st.button("Clear Data"):
            st.session_state.leads_df = None
            st.rerun()

    # --- MAIN PAGE UI ---
    st.title("🚀 Email Generator")

    col1, col2 = st.columns([1.2,1.8])
    with col1:
        st.markdown("#### 📎 Upload Documents")
        uploaded_files = st.file_uploader(
            "Upload Reference Files (TXT, PDF, DOCX, PPTX)",
            type=["txt", "pdf", "docx", "pptx"],
            accept_multiple_files=True,
            label_visibility="collapsed",
            key="blog_files"
        )
        
        file_context = ""
        if uploaded_files:
            for f in uploaded_files:
                extracted_text = extract_text_from_file(f) 
                file_context += f"--- Content from {f.name} ---\n{extracted_text}\n\n"

    with col2:
        st.markdown("#### 🔗 Reference URLs")
        reference_urls = st.text_area(
            "Add Reference URLs (comma-separated)",
            placeholder="https://example.com/page1, https://example.com/page2",
            height=70,
            label_visibility="collapsed",
            key="blog_urls"
        )
        url_list = [url.strip() for url in reference_urls.split(",") if url.strip()]


    st.markdown("### ✍️ Query & Instructions")
    query = st.text_area("Email Topic", key="email_topic", height=100)

    st.divider()

    # --- BULK GENERATION LOGIC ---
    if uploaded_file and st.session_state.leads_df is None:
        try:
            df = pd.read_csv(uploaded_file)
            st.session_state.leads_df = df
            st.success("CSV Loaded! Ready to generate.")
        except Exception as e:
            st.error(f"Error reading CSV: {e}")

    if st.session_state.leads_df is not None:
        df = st.session_state.leads_df
        
        common_params = {
            "query": query,
            "tone": tone,
            "word_limit": word_limit,
            "cta_choice": cta_choice,
            "reference_urls": url_list,
            "reference_file_content": file_context 
        }

        if "Generated Email" not in df.columns:
            if st.button("🚀 Generate Email", type="primary", use_container_width=True):
                if not query:
                    st.warning("Please enter an Email Topic before generating.")
                else:
                    with st.spinner("Generating emails... this may take a moment"):
                        st.session_state.leads_df = generate_bulk_emails(df, common_params)
                    st.rerun()
                
        else:
            st.subheader("Review Queue")
            new_df = df.drop(columns=['Generated Email'], errors='ignore')
            selection = st.dataframe(
                new_df,
                use_container_width=True,
                hide_index=True,
                on_select="rerun", 
                selection_mode="single-row",
                column_config={
                    "first_name": st.column_config.TextColumn("First Name"),
                    "last_name": st.column_config.TextColumn("Last Name"),
                    "org_name": st.column_config.TextColumn("Organization"),
                    "email": st.column_config.TextColumn("Email"),
                    "Status": st.column_config.SelectboxColumn(
                        "Status", 
                        options=["Draft", "Refined", "Approved", "Sent", "Failed"],
                        width="small"
                    )
                }
            )

            if selection.selection.rows:
                selected_index = selection.selection.rows[0]
                selected_row = df.iloc[selected_index]
                email_editor_dialog(selected_index, selected_row, common_params)